"""
ppt_writer.py
─────────────
Replaces placeholder charts and instruction boxes in PowerPoint slides.

Priority:
  1. If a native PowerPoint CHART shape is found overlapping or close to the instruction box,
     we update its categories and series data directly (using pptx.chart.data.CategoryChartData).
     This keeps the chart native, vector-sharp, style-consistent, and fully editable in PowerPoint.
  2. If no native CHART shape is found, we fall back to rendering a Matplotlib image,
     inserting it as a transparent PNG, and aligning it to the instruction shape.
  3. The original instruction box shape is deleted in both cases.
"""

import os
from typing import List, Dict, Any

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches
from rapidfuzz import fuzz, process

# Bounding box offset/distance in inches for pairing chart placeholders
PAIRING_DISTANCE_THRESHOLD = 3.5 


def _shapes_overlap(s1, s2) -> bool:
    """Check if two shape bounding boxes overlap in 2D space."""
    try:
        l1, t1, r1, b1 = s1.left, s1.top, s1.left + s1.width, s1.top + s1.height
        l2, t2, r2, b2 = s2.left, s2.top, s2.left + s2.width, s2.top + s2.height
        return not (r1 < l2 or r2 < l1 or b1 < t2 or b2 < t1)
    except Exception:
        return False


def _center_distance(s1, s2) -> float:
    """Calculate the distance in inches between the centers of two shapes."""
    try:
        c1_x = (s1.left + s1.width / 2) / 914400
        c1_y = (s1.top + s1.height / 2) / 914400
        c2_x = (s2.left + s2.width / 2) / 914400
        c2_y = (s2.top + s2.height / 2) / 914400
        return ((c1_x - c2_x)**2 + (c1_y - c2_y)**2)**0.5
    except Exception:
        return float('inf')


def _update_native_chart(chart, df: pd.DataFrame):
    """Update a native PowerPoint chart's categories and series values directly."""
    chart_data = CategoryChartData()
    
    # 1. Categories: Labels from DataFrame, excluding base rows (sample size counts)
    labels = []
    indices = []
    for idx, lbl in enumerate(df["Label"]):
        if "base:" not in str(lbl).lower():
            labels.append(str(lbl).strip())
            indices.append(idx)
            
    chart_data.categories = labels
    
    # 2. Match original chart series to excel columns
    orig_series = [s.name for s in chart.series]
    excel_cols = [c for c in df.columns if c != "Label"]
    
    for s_name in orig_series:
        best_col = None
        # Match series name to Excel columns
        best = process.extractOne(s_name, excel_cols, scorer=fuzz.ratio)
        if best and best[1] >= 60:
            best_col = best[0]
            
        if not best_col:
            # Fallback to Total or first available column
            best_col = "Total" if "Total" in excel_cols else excel_cols[0]
            
        vals = []
        for idx in indices:
            val = df[best_col].iloc[idx]
            try:
                v_num = float(val) if val is not None and not pd.isna(val) else 0.0
                # Scale proportions (between 0 and 1) to percentages
                if v_num > 0 and v_num <= 1.0:
                    v_num *= 100.0
                vals.append(v_num)
            except Exception:
                vals.append(0.0)
                
        chart_data.add_series(s_name, tuple(vals))
        
    chart.replace_data(chart_data)


def paste_charts_into_pptx(
    input_pptx_path: str,
    match_results: List[Dict[str, Any]],
    output_pptx_path: str,
) -> str:
    """
    Update native charts or insert Matplotlib fallback images on slide match results.
    """
    prs = Presentation(input_pptx_path)

    # Group matches by slide to process slide-by-slide
    slides_matches = {}
    for r in match_results:
        if r.get("match_confidence") != "unmatched":
            slide_num = r["slide_number"]
            if slide_num not in slides_matches:
                slides_matches[slide_num] = []
            slides_matches[slide_num].append(r)

    for slide_num, matches in slides_matches.items():
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            continue

        slide = prs.slides[slide_idx]
        
        # Keep track of shapes we will delete from this slide
        shapes_to_delete = []
        fallback_charts_to_insert = [] # tuple of (image_path, left, top, width, height)

        for result in matches:
            # Find the instruction shape by shape_id
            instr_shape = None
            for shape in slide.shapes:
                if shape.shape_id == result["shape_id"]:
                    instr_shape = shape
                    break

            if not instr_shape:
                continue

            # Look for paired native CHART shapes on the slide
            paired_chart_shape = None
            min_dist = float('inf')

            for shape in slide.shapes:
                if shape.shape_type == 3:  # CHART
                    overlap = _shapes_overlap(instr_shape, shape)
                    dist = _center_distance(instr_shape, shape)
                    
                    if overlap or dist < Inches(PAIRING_DISTANCE_THRESHOLD):
                        if dist < min_dist:
                            min_dist = dist
                            paired_chart_shape = shape

            data_dict = result.get("matched_table_data")
            
            # --- Strategy 1: Update native PowerPoint chart if found ---
            if paired_chart_shape and data_dict:
                try:
                    df = pd.DataFrame(data_dict)
                    _update_native_chart(paired_chart_shape.chart, df)
                    shapes_to_delete.append(instr_shape) # delete instruction box
                    continue # successfully updated native, skip image fallback
                except Exception as e:
                    print(f"Failed to update native chart on Slide {slide_num}: {e}. Falling back to image...")

            # --- Strategy 2: Fallback to Matplotlib transparent image ---
            chart_path = result.get("chart_path")
            if chart_path and os.path.isfile(chart_path):
                target_left = instr_shape.left
                target_top = instr_shape.top
                target_width = instr_shape.width
                target_height = instr_shape.height
                
                # If there's a chart placeholder, align image to it
                if paired_chart_shape:
                    target_left = paired_chart_shape.left
                    target_top = paired_chart_shape.top
                    target_width = paired_chart_shape.width
                    target_height = paired_chart_shape.height
                    shapes_to_delete.append(paired_chart_shape)
                    
                shapes_to_delete.append(instr_shape)
                fallback_charts_to_insert.append((chart_path, target_left, target_top, target_width, target_height))

        # Insert any fallback images
        for chart_path, left, top, width, height in fallback_charts_to_insert:
            try:
                slide.shapes.add_picture(chart_path, left, top, width, height)
            except Exception as e:
                print(f"Error inserting chart image on Slide {slide_num}: {e}")

        # Delete instruction shapes and replaced placeholders
        for shape in shapes_to_delete:
            try:
                slide.shapes._spTree.remove(shape._element)
            except Exception as e:
                print(f"Error deleting shape ID {shape.shape_id} on Slide {slide_num}: {e}")

    prs.save(output_pptx_path)
    return output_pptx_path
