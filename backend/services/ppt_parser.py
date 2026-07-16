"""
ppt_parser.py
─────────────
Extracts ALL text-bearing shapes from a .pptx file.

Handles:
  • Standard text boxes
  • AutoShapes (callouts, arrows, rounded rectangles, etc.)
  • Grouped shapes  (recursive descent)
  • Table cells inside shapes
  • Freeform / SmartArt-like shapes

Deduplication:
  A SHA-256 hash of (slide_number + shape_id + raw_text) is computed
  for every candidate. Entries whose hash already exists in the seen-set
  are silently dropped so the same instruction can never appear twice.

Classification:
  "instruction" → text contains {Data}, Table, Q-codes (Q1, S1 …), Column, Total
  "symbol"      → text is only/mainly a directional or special character
  "unknown"     → anything else
"""

import hashlib
import re
from typing import List, Dict, Any

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── Patterns ──────────────────────────────────────────────────────────────────
SYMBOL_PATTERN = re.compile(
    r"^[\s↑↓→←▲▼◀▶⬆⬇⬅➡★☆✓✗✔✘%#@&\*\+\-=<>!?~^|•·…\u2019\u201c\u201d]+$"
)
INSTRUCTION_PATTERN = re.compile(
    r"(\{[Dd]ata\}|[Tt]able[\s_]?\d+|[QqSs]\d+[a-zA-Z]?|[Cc]olumn|[Tt]otal|[Cc]ross\s*[Tt]ab|"
    r"[Gg]ender|[Aa]ge|[Mm]arital|[Ee]duc|[Cc]oncern|[Rr]ank|[Ss]tatement|[Rr]eason)",
    re.IGNORECASE,
)


def _classify(text: str) -> str:
    stripped = text.strip()
    if not stripped:
        return None
    if SYMBOL_PATTERN.match(stripped):
        return "symbol"
    if INSTRUCTION_PATTERN.search(stripped):
        return "instruction"
    return "unknown"


def _make_hash(slide_number: int, shape_id: int, raw_text: str) -> str:
    payload = f"{slide_number}|{shape_id}|{raw_text}"
    return hashlib.sha256(payload.encode()).hexdigest()


def _shape_bounds(shape) -> Dict[str, float]:
    """Return position/size in centimetres (float, 2dp)."""
    def emu_to_cm(emu):
        return round(emu / 914400 * 2.54, 4)
    try:
        return {
            "left": emu_to_cm(shape.left or 0),
            "top": emu_to_cm(shape.top or 0),
            "width": emu_to_cm(shape.width or 0),
            "height": emu_to_cm(shape.height or 0),
        }
    except Exception:
        return {"left": 0, "top": 0, "width": 0, "height": 0}


def _extract_from_shape(
    shape,
    slide_number: int,
    seen: set,
    results: list,
    counter: list,  # mutable int holder
):
    """Recursively extract text from a single shape."""
    shape_type = shape.shape_type

    # ── Group: recurse into children ──────────────────────────────────────
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _extract_from_shape(child, slide_number, seen, results, counter)
        return

    # ── Table: iterate cells ───────────────────────────────────────────────
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text = cell.text_frame.text.strip() if cell.text_frame else ""
                if text:
                    _record(text, shape, slide_number, seen, results, counter)
        return

    # ── Text frame ─────────────────────────────────────────────────────────
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            _record(text, shape, slide_number, seen, results, counter)


def _record(text, shape, slide_number, seen, results, counter):
    kind = _classify(text)
    if kind is None:
        return

    h = _make_hash(slide_number, shape.shape_id, text)
    if h in seen:
        return  # ← deduplication: skip identical instruction
    seen.add(h)

    counter[0] += 1
    bounds = _shape_bounds(shape)
    results.append({
        "id": counter[0],
        "slide_number": slide_number,
        "shape_id": shape.shape_id,
        "shape_name": shape.name or f"Shape_{shape.shape_id}",
        "raw_text": text,
        "type": kind,
        "hash": h,
        **bounds,
    })


def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
    """
    Main entry point.
    Returns a list of unique instruction/symbol dicts from the .pptx file.
    """
    prs = Presentation(file_path)
    results: list = []
    seen: set = set()
    counter = [0]  # mutable counter passed by reference

    for slide_index, slide in enumerate(prs.slides):
        slide_number = slide_index + 1
        for shape in slide.shapes:
            _extract_from_shape(shape, slide_number, seen, results, counter)

    return results
